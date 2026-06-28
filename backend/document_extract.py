"""
document_extract.py — Trích text thuần từ file Word/Excel/PowerPoint.

BỐI CẢNH: UploadPage (App.jsx) cho phép chọn .docx/.xlsx/.pptx (input accept +
chip định dạng hiển thị), nhưng trước đây backend CHỈ xử lý PDF — file khác
được gửi lên thì bị từ chối với lỗi "Chỉ chấp nhận file PDF", hiển thị ở FE
dưới dạng thông báo lỗi server chung (gây hiểu nhầm là lỗi server, không phải
lỗi loại file). Module này lấp đúng khoảng trống đó cho 3 định dạng có text
sẵn trong file (không phải ảnh scan) — ẢNH (.png/.jpg) KHÔNG nằm trong phạm vi
module này vì cần OCR thật (chưa có engine OCR local hay key SmartReader);
proposal đã tự ghi rõ OCR "để dành giai đoạn 2" — giữ đúng định hướng đó.

Mỗi hàm trả (text: str, warning: str | None) — warning khi file rỗng hoặc lỗi
đọc một phần, để main.py quyết định có nên tiếp tục pipeline hay báo lỗi rõ
cho người dùng.
"""
from typing import Tuple, Optional
import io

from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl


def extract_docx(file_bytes: bytes) -> Tuple[str, Optional[str]]:
    """Trích text từ file Word (.docx). Lấy cả đoạn văn và bảng (bệnh án Word
    thường có bảng xét nghiệm/thuốc dạng table, không chỉ đoạn văn)."""
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))
    except Exception as e:
        return "", f"Không đọc được file Word: {e}"

    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())

    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    text = "\n".join(parts)
    if not text.strip():
        return "", "File Word không có nội dung văn bản đọc được (có thể là file scan/ảnh nhúng)."
    return text, None


def extract_xlsx(file_bytes: bytes) -> Tuple[str, Optional[str]]:
    """Trích text từ Excel (.xlsx). Mỗi sheet -> tiêu đề sheet + các dòng, ô
    nối bằng tab để giữ cấu trúc bảng tương đối (giúp Claude đọc đúng cột)."""
    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    except Exception as e:
        return "", f"Không đọc được file Excel: {e}"

    parts = []
    for sheet in wb.worksheets:
        sheet_lines = [f"--- Sheet: {sheet.title} ---"]
        has_content = False
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip() != ""]
            if cells:
                has_content = True
                sheet_lines.append("\t".join(cells))
        if has_content:
            parts.append("\n".join(sheet_lines))

    text = "\n\n".join(parts)
    if not text.strip():
        return "", "File Excel không có dữ liệu (tất cả sheet đều rỗng)."
    return text, None


def extract_pptx(file_bytes: bytes) -> Tuple[str, Optional[str]]:
    """Trích text từ PowerPoint (.pptx). Mỗi slide -> số thứ tự + text từng
    shape (bao gồm cả bảng nếu có)."""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        return "", f"Không đọc được file PowerPoint: {e}"

    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = [f"--- Slide {i} ---"]
        has_content = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                has_content = True
                slide_lines.append(shape.text_frame.text.strip())
            elif shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        has_content = True
                        slide_lines.append(" | ".join(cells))
        if has_content:
            parts.append("\n".join(slide_lines))

    text = "\n\n".join(parts)
    if not text.strip():
        return "", "File PowerPoint không có nội dung văn bản đọc được."
    return text, None


# Map đuôi file -> hàm trích text, để main.py gọi thống nhất không cần if/elif dài.
EXTRACTORS = {
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".pptx": extract_pptx,
}

SUPPORTED_EXTENSIONS = set(EXTRACTORS.keys()) | {".pdf"}
# Ảnh CHƯA hỗ trợ — liệt kê riêng để main.py trả lỗi đúng nghĩa "chưa hỗ trợ"
# thay vì lỗi chung, và để FE có thể đồng bộ lại accept/chip hiển thị sau.
UNSUPPORTED_BUT_LISTED_IN_UI = {".png", ".jpg", ".jpeg", ".doc", ".xls", ".ppt"}
# Lưu ý: .doc/.xls/.ppt (định dạng cũ, không phải Open XML .docx/.xlsx/.pptx)
# KHÔNG được hỗ trợ bởi python-docx/openpyxl/python-pptx — cần thư viện khác
# (vd antiword, xlrd cũ) nếu muốn làm tiếp; tạm xếp vào nhóm "chưa hỗ trợ".


def extract_from_filename(filename: str, file_bytes: bytes) -> Tuple[str, Optional[str], bool]:
    """
    Điểm vào duy nhất: nhận tên file + bytes, tự chọn extractor theo đuôi.
    Trả (text, warning, supported) — supported=False nghĩa là loại file này
    chưa có extractor (ảnh, .doc/.xls/.ppt cũ).
    """
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    extractor = EXTRACTORS.get(ext)
    if extractor is None:
        return "", None, False
    text, warning = extractor(file_bytes)
    return text, warning, True
